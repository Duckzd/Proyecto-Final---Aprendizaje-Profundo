"""
Genera la presentación EDITABLE (elementos nativos de PowerPoint).

Ejecutar desde la raíz del proyecto:
    python ppt/build_editable.py

Sólo son imágenes los dibujos del ninja, los esqueletos y el GIF
(están en ppt/native/ y ppt/img/). Todo el texto, cajas, tablas y
flechas son objetos nativos y editables.
"""
import os, sys
sys.path.insert(0, os.getcwd())
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image
import config, character, kinematics

N = "ppt/native/"
BG=RGBColor(0xF6,0xF3,0xEE); INK=RGBColor(0x1B,0x24,0x30); SLATE=RGBColor(0x3A,0x4B,0x5C)
BRONZE=RGBColor(0xA8,0x76,0x3E); SAGE=RGBColor(0x5E,0x8B,0x72); CLAY=RGBColor(0xA8,0x54,0x48)
MUTED=RGBColor(0x8A,0x8F,0x98); LINE=RGBColor(0xD8,0xD2,0xC8); WHITE=RGBColor(0xFF,0xFF,0xFF)
CARD=RGBColor(0xFF,0xFF,0xFF); SLATE2=RGBColor(0x4A,0x62,0x74)
GREY1=RGBColor(0x7A,0x8B,0x96); GREY2=RGBColor(0x9A,0xA5,0xAC)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    return s

def txt(s,l,t,w,h,text,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(str(text).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name="Calibri"
    return tb

def box(s,l,t,w,h,fill=None,line=None,rounded=True,lw=1.25,radius=0.09):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                          Inches(l),Inches(t),Inches(w),Inches(h))
    sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    if rounded:
        try: sh.adjustments[0]=radius
        except Exception: pass
    return sh

def boxtext(sh,text,size,color,bold=True,align=PP_ALIGN.CENTER):
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,ln in enumerate(str(text).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name="Calibri"

def pic(s,path,l,t,w=None,h=None):
    iw,ih=Image.open(path).size; ar=iw/ih
    if w and not h: h=w/ar
    if h and not w: w=h*ar
    return s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))

def pic_c(s,path,cx,t,h):
    iw,ih=Image.open(path).size; w=h*iw/ih
    s.shapes.add_picture(path,Inches(cx-w/2),Inches(t),Inches(w),Inches(h)); return w

def arrow(s,x1,y1,x2,y2,color=MUTED,w=2.0,dash=False):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w)
    ln=c.line._get_or_add_ln()
    if dash: ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))
    return c

def hline(s,x1,x2,y,color=LINE,w=1.0):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y),Inches(x2),Inches(y))
    c.line.color.rgb=color; c.line.width=Pt(w)

def header(s,sec,title,sub=None):
    txt(s,0.55,0.26,11,0.34,sec.upper(),12,BRONZE,True)
    txt(s,0.55,0.55,12.2,0.72,title,26,INK,True)
    if sub: txt(s,0.57,1.22,12.2,0.4,sub,13,MUTED,italic=True)
    hline(s,0.6,12.73,1.66)

def banner(s,text,color=INK):
    w=min(12.5,0.7+len(text)*0.089); l=(13.333-w)/2
    sh=box(s,l,6.74,w,0.5,fill=color,radius=0.5); boxtext(sh,text,12.5,WHITE,True)

def clear_tbl_style(tbl):
    pr=tbl._tbl.tblPr
    for e in pr.findall(qn('a:tableStyleId')): pr.remove(e)
    sid=pr.makeelement(qn('a:tableStyleId'),{}); sid.text="{2D5ABB26-0587-4C30-8999-92F81FD0307C}"; pr.append(sid)

def table(s,l,t,w,h,data,col_w,fonts=9):
    rows,cols=len(data),len(data[0])
    g=s.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(h)); tbl=g.table
    tbl.first_row=False; tbl.horz_banding=False; clear_tbl_style(tbl)
    for j,cw in enumerate(col_w): tbl.columns[j].width=Inches(cw)
    for i in range(rows):
        tbl.rows[i].height=Inches(h/rows)
        for j in range(cols):
            c=tbl.cell(i,j)
            c.margin_top=Pt(1); c.margin_bottom=Pt(1); c.margin_left=Pt(6); c.margin_right=Pt(6)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE; c.fill.solid(); c.fill.fore_color.rgb=CARD
            p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
            r=p.add_run(); r.text=str(data[i][j]); r.font.size=Pt(fonts); r.font.name="Calibri"; r.font.color.rgb=SLATE
    return tbl

def cell_color(tbl,i,fill,fg):
    for j in range(len(tbl.columns)):
        c=tbl.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=fill
        for p in c.text_frame.paragraphs:
            for r in p.runs: r.font.color.rgb=fg; r.font.bold=True

STO=config.STATE_ORDER
def unit(i): return "°" if (i>=2 and i not in (3,4,16)) else (" px" if i<2 else "")

KICK=np.load(N+"kick_mat.npy"); F=22

# ============================================== 1 PORTADA
s=slide()
txt(s,0.8,0.85,11.7,1.0,"Animar un ninja 2D",40,INK,True,PP_ALIGN.CENTER)
txt(s,0.8,1.78,11.7,0.6,"con una red neuronal recurrente",24,BRONZE,False,PP_ALIGN.CENTER)
hline(s,4.9,8.43,2.5,LINE,1.5)
x=3.15
for f in ["p_idle.png","p_slash.png","p_kick.png"]:
    pic_c(s,N+f,x+1.15,2.75,3.1); x+=2.55
txt(s,0.8,6.45,11.7,0.4,"Proyecto Final  ·  Aprendizaje Profundo",14,MUTED,False,PP_ALIGN.CENTER)

# ============================================== 2 VISIÓN GENERAL
s=slide(); header(s,"Visión general","En una idea: enseñar a una red a mover un ninja")
txt(s,0.6,1.85,12,0.4,"Animar un personaje a mano, fotograma a fotograma, es lento y tedioso.",14,SLATE)
txt(s,0.6,2.30,12,0.4,"La meta: que una red neuronal aprenda a generar el movimiento ella sola.",14,INK,True)
pic_c(s,N+"p_kick.png",1.55,3.05,1.95)
by=3.35; bh=1.2
for l,c,ti,su in [(2.95,SLATE,"DATOS","cada pose = 17 ángulos"),
                  (5.35,BRONZE,"RED NEURONAL","GRU recurrente"),
                  (7.75,SAGE,"DIBUJO","cinemática directa")]:
    b=box(s,l,by,2.0,bh,fill=c); boxtext(b,ti+"\n"+su,12.5,WHITE,True)
pic_c(s,N+"p_slash.png",11.5,3.05,1.95)
for x1,x2 in [(2.55,2.9),(4.95,5.30),(7.35,7.70),(9.75,10.35)]:
    arrow(s,x1,by+bh/2,x2,by+bh/2)
for cx,tt in [(1.55,"movimientos\nde ejemplo"),(3.95,"nada de coordenadas:\nÁNGULOS"),
              (6.35,"aprende a predecir\nla pose siguiente"),(8.75,"de los ángulos\na la imagen")]:
    txt(s,cx-1.1,by+bh+0.05,2.2,0.6,tt,9.5,MUTED,italic=True,align=PP_ALIGN.CENTER)
txt(s,10.5,by+bh+0.05,2.0,0.6,"¡animación\ngenerada!",9.5,SAGE,True,PP_ALIGN.CENTER)
txt(s,0.6,5.35,12,0.34,"EL RECORRIDO DE HOY",11,BRONZE,True)
rx=0.6; rw=2.42
for num,tt,c in [("I","El problema",CLAY),("II","Los datos",SLATE),("III","Del dato al dibujo",BRONZE),
                 ("IV","El modelo",SAGE),("V","Los resultados",INK)]:
    box(s,rx,5.72,rw-0.12,0.62,fill=CARD,line=c,lw=1.75)
    cc=box(s,rx+0.14,5.86,0.35,0.35,fill=c,radius=0.5); boxtext(cc,num,11,WHITE,True)
    txt(s,rx+0.58,5.72,rw-0.7,0.62,tt,11.5,INK,True,anchor=MSO_ANCHOR.MIDDLE)
    rx+=rw
banner(s,"Todo el sistema se construye sobre una idea: representar el cuerpo con ÁNGULOS, no con coordenadas")

# ============================================== 3 EL PROBLEMA
s=slide(); header(s,"I · el problema","La idea que falló: predecir las coordenadas (x, y)",
                  "basta un error pequeño en cada punto para que los huesos cambien de longitud")
txt(s,0.6,1.9,6.0,0.4,"Prediciendo ÁNGULOS",15,SAGE,True,PP_ALIGN.CENTER)
txt(s,6.7,1.9,6.0,0.4,"Prediciendo COORDENADAS (x, y)",15,CLAY,True,PP_ALIGN.CENTER)
pic_c(s,N+"prob_angles.png",3.6,2.35,3.7)
pic_c(s,N+"prob_coords.png",9.7,2.35,3.7)
txt(s,0.6,6.15,6.0,0.4,"antebrazo 65 px      pantorrilla 90 px",13,SAGE,True,PP_ALIGN.CENTER)
txt(s,6.7,6.15,6.0,0.4,"antebrazo 86 px      pantorrilla 87 px",13,CLAY,True,PP_ALIGN.CENTER)
banner(s,"Solución: predecir los ÁNGULOS. Con cinemática directa las longitudes quedan fijas por construcción.",SAGE)

# ============================================== 4 UNA POSE
s=slide(); header(s,"II · los datos","Una POSE son 17 números",
                  "posición de la cadera, velocidades y diez ángulos que describen el cuerpo entero")
data=[[STO[i], f"{KICK[F,i]:.1f}{unit(i)}"] for i in range(17)]
tbl=table(s,0.6,1.9,3.5,4.6,data,[2.2,1.3],fonts=9.5)
for i,c in {5:SAGE,11:SLATE,13:CLAY,14:BRONZE}.items():
    cell_color(tbl,i,RGBColor(0xF0,0xEC,0xE4),c)
IL,IT,IH=4.6,1.85,4.55; CROP=(120,235,690,800)
iw=CROP[2]-CROP[0]; ih=CROP[3]-CROP[1]; IW=IH*iw/ih
pic(s,N+"pose_kick.png",IL,IT,w=IW,h=IH)
po=kinematics.forward_kinematics(character.state_from_vector(KICK[F],'k')); P=po.positions
def j2s(n):
    x,y=P[n]; return IL+(x-CROP[0])/iw*IW, IT+((config.WINDOW_HEIGHT-y)-CROP[1])/ih*IH
mid=lambda a,b: ((j2s(a)[0]+j2s(b)[0])/2,(j2s(a)[1]+j2s(b)[1])/2)
for tt,c,bl,bt,(jx,jy) in [("torso_angle = 9.9°",SAGE,3.9,2.05,mid('pelvis','chest')),
                           ("right_hip = 103.9°",CLAY,9.95,2.05,mid('pelvis','r_knee')),
                           ("right_knee = 8.3°",BRONZE,10.25,4.35,j2s('r_ankle')),
                           ("left_knee = 32.9°",SLATE,3.7,5.15,j2s('l_knee'))]:
    b=box(s,bl,bt,2.75,0.5,fill=c,radius=0.25); boxtext(b,tt,11.5,WHITE,True)
    arrow(s,(bl+2.75 if bl<7 else bl),bt+0.25,jx,jy,c,1.8)
banner(s,"cada ÁNGULO se mide respecto al hueso anterior  ·  la cadera (root) da la posición: (400, 248)")

# ============================================== 5 ¿RESPECTO A QUÉ GIRA?
s=slide(); header(s,"II · los datos","¿Respecto a qué gira cada ángulo?",
                  "cada ángulo se mide desde la prolongación del hueso PADRE, no desde el suelo")
pic(s,N+"s05_angles.png",0.5,1.8,h=4.7)
lx=8.55; ly=2.2
txt(s,lx,ly-0.4,4.2,0.35,"CÓMO LEER LA FIGURA",11,BRONZE,True)
for c,ti,de in [(MUTED,"eje GLOBAL","referencia fija (horizontal / vertical), igual para todo el cuerpo"),
                (SLATE2,"eje RELATIVO","la prolongación del hueso padre: el '0°' de esa articulación"),
                (BRONZE,"arco = el ÁNGULO","cuánto se ha girado respecto a ese eje relativo")]:
    box(s,lx,ly+0.03,0.32,0.32,fill=c,radius=0.3)
    txt(s,lx+0.45,ly-0.05,3.9,0.35,ti,13,c,True)
    txt(s,lx+0.45,ly+0.30,3.95,0.62,de,10.5,SLATE)
    ly+=1.0
txt(s,lx,ly+0.1,4.2,0.9,"Los números (espalda 25°, codo 90°, rodilla 80°) son justo los valores del estado.",11,INK,True)
banner(s,"Por eso son ángulos RELATIVOS: al girar un hueso, todo lo que cuelga de él le acompaña.",BRONZE)

# ============================================== 6 UNA SECUENCIA
s=slide(); header(s,"II · los datos","Una SECUENCIA son 40 poses seguidas",
                  "los mismos 17 números, cambiando poco a poco en cada fotograma")
xs=0.55; step=2.06
for k,f in enumerate([0,8,16,22,30,39]):
    cx=xs+step*k+0.95; hl=(f==22)
    txt(s,cx-1.0,1.85,2.0,0.5,f"fotograma {f}"+("  ¡impacto!" if hl else ""),10.5,(CLAY if hl else SLATE),True,PP_ALIGN.CENTER)
    pic_c(s,N+f"seq{f}.png",cx,2.35,2.7)
    b=box(s,cx-0.95,5.25,1.9,1.15,fill=CARD,line=(CLAY if hl else LINE),lw=(2.2 if hl else 1.2))
    boxtext(b,f"right_hip\n{KICK[f,13]:.1f}°\nright_knee\n{KICK[f,14]:.1f}°",9.5,(CLAY if hl else SLATE),True)
banner(s,"Al reproducirlas a 30 fotogramas por segundo, el ojo las une y ve UNA PATADA")

# ============================================== 7 EL DATASET
s=slide(); header(s,"II · los datos","El DATASET son 600 secuencias como ésa",
                  "10 movimientos distintos y, de cada uno, 60 variantes aleatorias: ninguna igual a otra")
movs=["idle","walk","run","jump","dash","roll","punch","kick","sword_slash","sword_combo"]
xs=0.35; stp=1.29
for k,mv in enumerate(movs):
    cx=xs+stp*k+0.55
    txt(s,cx-0.85,1.85,1.7,0.35,mv,9.5,(CLAY if mv=="kick" else SLATE),True,PP_ALIGN.CENTER)
    pic_c(s,N+f"ds_{mv}.png",cx,2.25,2.35)
txt(s,0.6,4.75,12.1,0.5,"10 movimientos   ×   60 variantes cada uno   =   600 secuencias",15,INK,True,PP_ALIGN.CENTER)
bx=2.9
for num,tt,c in [("600","secuencias\n(ejecuciones)",SLATE),("40","fotogramas\n(instantes)",BRONZE),("17","números\n(la pose)",SAGE)]:
    b=box(s,bx,5.45,2.1,1.05,fill=c); boxtext(b,num+"\n"+tt,13,WHITE,True)
    if bx<6.5: txt(s,bx+2.1,5.45,0.5,1.05,"×",18,INK,True,PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bx+=2.6
banner(s,"cada casilla es UN número:   data[secuencia, fotograma, variable]")

# ============================================== 8 FK
s=slide(); header(s,"III · del dato al dibujo","Paso 1: de los ángulos a los PUNTOS",
                  "cada hueso hereda la orientación del padre; su ángulo GLOBAL es el que entra en la fórmula")
pic(s,N+"fk_chain.png",0.3,1.8,h=4.7)
b=box(s,6.2,1.9,6.55,0.95,fill=CARD,line=INK,lw=1.6)
boxtext(b,"La regla (trigonometría)\npunto hijo  =  punto padre  +  longitud × (cos ángulo, sen ángulo)",12,INK,True)
ty=3.05
for num,c,t1,t2 in [("1",INK,"Empezamos en la cadera","cadera = (400, 248)"),
                    ("2",CLAY,"El muslo mide 95 px y apunta a 5.9°","rodilla = (400,248) + 95×(cos 5.9°, sen 5.9°) = (494, 258)"),
                    ("3",BRONZE,"La pantorrilla mide 90 px y apunta a −2.4°","tobillo = (494,258) + 90×(cos −2.4°, sen −2.4°) = (584, 254)")]:
    b=box(s,6.2,ty,6.55,1.05,fill=CARD,line=c,lw=1.6)
    cc=box(s,6.4,ty+0.33,0.4,0.4,fill=c,radius=0.5); boxtext(cc,num,13,WHITE,True)
    txt(s,6.95,ty+0.08,5.7,0.4,t1,11.5,c,True)
    txt(s,6.95,ty+0.5,5.7,0.5,t2,10.5,SLATE)
    ty+=1.2
banner(s,"…y así con los 17 nodos. Las longitudes NUNCA cambian: sólo los ángulos.")

# ============================================== 9 RENDER
s=slide(); header(s,"III · del dato al dibujo","Paso 2: de los puntos al NINJA",
                  "cada pieza PNG se gira el ángulo de su hueso y se pega por su PIVOTE en la articulación")
for k,(ti,su,fn,c) in enumerate([("1. los PUNTOS","(lo que da el paso anterior)","render_points.png",SLATE),
                                 ("2. los HUESOS","(unir punto con punto)","render_bones.png",SLATE),
                                 ("3. el NINJA","(pegar las piezas PNG)","pose_kick.png",BRONZE)]):
    cx=2.4+k*3.9
    txt(s,cx-1.7,1.85,3.4,0.6,ti+"\n"+su,12,c,True,PP_ALIGN.CENTER)
    pic_c(s,N+fn,cx,2.75 if k==2 else 2.85,2.85 if k==2 else 2.7)
    if k<2: arrow(s,cx+1.55,4.1,cx+1.95,4.1,MUTED,2.2)
px=0.65; slot=1.4
for p in ["head","torso","upper_arm","forearm","hand","thigh","calf","foot","sword"]:
    iw2,ih2=Image.open("assets/%s.png"%p).size; ar=iw2/ih2
    h=0.5; w=h*ar
    if w>1.2: w=1.2; h=w/ar
    pic(s,"assets/%s.png"%p,px+(slot-w)/2,6.0-h/2,w=w,h=h)
    txt(s,px,6.28,slot,0.3,p,8,MUTED,align=PP_ALIGN.CENTER); px+=slot
banner(s,"Como el pivote de cada pieza cae sobre la articulación, nunca quedan huecos entre partes")

# ============================================== 10 VENTANAS (entrenamiento)
s=slide(); header(s,"IV · el modelo","Cómo se le enseña: de 600 secuencias a 12.240 ejemplos",
                  "cada secuencia se trocea en ventanas; la red NUNCA ve una secuencia entera")
txt(s,2.40,1.72,3.5,0.30,"ENTRADA (las últimas poses)",10.5,SLATE,True,PP_ALIGN.CENTER)
txt(s,7.65,1.72,2.0,0.30,"OBJETIVO",10.5,BRONZE,True,PP_ALIGN.CENTER)
frames=[0,3,6,9,12,15,18,21]; xs=0.55; stp=1.5
for row,shift in enumerate([0,1]):
    yb=1.98+row*1.04
    for j,f in enumerate(frames):
        cx=xs+stp*j+0.6
        pic_c(s,N+f"win{f}.png",cx,yb,0.92)
        if shift<=j<shift+5: box(s,cx-0.60,yb-0.05,1.20,1.02,fill=None,line=SLATE,lw=2.0)
        elif j==shift+5:     box(s,cx-0.60,yb-0.05,1.20,1.02,fill=None,line=BRONZE,lw=2.2)
    txt(s,12.30,yb+0.28,1.0,0.35,f"ejemplo {row+1}",10,MUTED,italic=True)
txt(s,0.6,4.04,12.1,0.34,"la ventana se desliza UN fotograma y ya tienes el siguiente ejemplo        "
                          "(aquí se ven 5 poses por espacio; la ventana real es de 16)",
    10.5,MUTED,italic=True,align=PP_ALIGN.CENTER)
py=4.50; ph=1.20; pw=2.15
stages=[("DATASET","(600, 40, 17)","600 secuencias",SLATE),
        ("SE PARTE","510 train · 90 val","por secuencias",SLATE2),
        ("SE TROCEA","24 ventanas","por cada secuencia",BRONZE),
        ("EJEMPLOS","X (12.240, 16, 17)","y (12.240, 17)",BRONZE),
        ("LOTES","(256, 16, 17)","barajados, a la red",SAGE)]
px=0.62
for ti,l2,l3,c in stages:
    b=box(s,px,py,pw,ph,fill=c); boxtext(b,ti+"\n"+l2+"\n"+l3,10,WHITE,True); px+=pw+0.35
for k in range(4):
    a=0.62+(k+1)*pw+k*0.35; arrow(s,a+0.04,py+ph/2,a+0.31,py+ph/2)
txt(s,0.6,5.83,12.1,0.36,"cada lote mezcla ventanas de secuencias y movimientos distintos: la red aprende una regla general, no el orden del dataset",
    10.5,SLATE,italic=True,align=PP_ALIGN.CENTER)
banner(s,"Entrenar = 12.240 ejemplos sueltos y barajados   ·   Generar = un bucle, una pose tras otra")

# ============================================== 11 ARQUITECTURA
s=slide(); header(s,"IV · el modelo","La arquitectura: entra una ventana, sale la pose siguiente",
                  "dos capas GRU resumen el movimiento reciente y dos capas densas producen los 17 números")
Y=3.05; H=1.5
for tt,c,l,w,sub in [("ENTRADA\n16 poses × 17",SLATE,0.55,1.75,"las últimas 16\nposes del ninja"),
                     ("GRU\ncapa 1 · 64",SLATE2,2.45,1.55,"lee la secuencia\ny recuerda el ritmo"),
                     ("GRU\ncapa 2 · 64",SLATE2,4.15,1.55,"refina ese\nresumen temporal"),
                     ("último\npaso",GREY1,5.85,1.35,"resume los 16\nen 1 vector"),
                     ("Dense 64\n+ ReLU",GREY2,7.35,1.25,"combina\nlas señales"),
                     ("Dense\n17",SAGE,8.75,1.15,"una salida\npor variable")]:
    b=box(s,l,Y,w,H,fill=c); boxtext(b,tt,12,WHITE,True)
    txt(s,l-0.1,Y+H+0.05,w+0.2,0.55,sub,9.2,MUTED,italic=True,align=PP_ALIGN.CENTER)
for x in [2.32,4.02,5.72,7.24,8.64]: arrow(s,x,Y+H/2,x+0.11,Y+H/2)
arrow(s,9.92,Y+H/2,10.35,Y+H/2,SAGE,2.4)
txt(s,10.4,Y+0.15,2.6,H,"la POSE\nSIGUIENTE\n(17 números)",13,SAGE,True,PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
fy=5.65
arrow(s,11.55,Y+H,11.55,fy,CLAY,2.2,dash=True)
hline(s,1.45,11.55,fy,CLAY,2.2)
arrow(s,1.45,fy,1.45,Y+H,CLAY,2.2,dash=True)
txt(s,3.0,fy-0.02,7.5,0.4,"AUTOREGRESIVO: la pose predicha vuelve a entrar",12,CLAY,True,PP_ALIGN.CENTER)
txt(s,0.6,6.05,12.1,0.35,"46.161 parámetros   ·   pérdida MSE   ·   optimizador Adam   ·   15 épocas",12,MUTED,False,PP_ALIGN.CENTER)
banner(s,"Repitiendo el paso se construye la animación completa, fotograma a fotograma")

# ============================================== 12 SALIDAS
s=slide(); header(s,"V · los resultados","¿Qué devuelve la red y cómo se interpreta?",
                  "no devuelve una imagen: devuelve 17 números que después convertimos en dibujo")
pred=np.load(N+"pred_vec.npy"); mean,std=np.load(N+"norm.npy"); pn=(pred-mean)/std
show=[0,1,5,13,14,16]
txt(s,0.6,1.9,3.6,0.55,"1 · lo que sale de la red\n(normalizado)",11.5,SLATE,True,PP_ALIGN.CENTER)
table(s,0.6,2.55,3.6,3.4,[[STO[i],f"{pn[i]:.2f}"] for i in show],[2.3,1.3],10)
txt(s,4.5,1.9,3.6,0.55,"2 · al deshacer la normalización\n(grados y píxeles)",11.5,BRONZE,True,PP_ALIGN.CENTER)
table(s,4.5,2.55,3.6,3.4,[[STO[i],f"{pred[i]:.1f}{unit(i)}"] for i in show],[2.3,1.3],10)
txt(s,8.4,1.9,4.4,0.55,"3 · se dibuja con cinemática directa\n(el fotograma predicho)",11.5,SAGE,True,PP_ALIGN.CENTER)
pic_c(s,N+"pred_pose.png",10.6,2.5,3.6)
arrow(s,4.25,4.2,4.45,4.2); arrow(s,8.15,4.2,8.35,4.2)
banner(s,"Para ver la animación: repetir estos 3 pasos por fotograma  ·  notebook, Sección 11 (reproductor)")

# ============================================== 13 RESULTADOS
s=slide(); header(s,"V · los resultados","La red reproduce el movimiento",
                  "se le dan 16 fotogramas de semilla y genera el resto ella sola")
idxs=[16,20,24,28,32,36,39]
for row,(pref,tt,c) in enumerate([("ro","ORIGINAL · generador matemático",SLATE),
                                  ("rg","GENERADO POR LA RED · autoregresivo",SAGE)]):
    yb=2.35+row*1.85
    txt(s,0.6,yb-0.4,7,0.3,tt,11.5,c,True)
    for k,f in enumerate(idxs): pic_c(s,N+f"{pref}{f}.png",0.9+k*1.15,yb,1.55)
tbl=table(s,9.5,2.35,3.3,2.0,[["modelo","MAE","RMSE"],["RNN","1.93","3.58"],["LSTM","1.73","3.11"],["GRU","1.78","3.30"]],[1.5,0.9,0.9],11)
cell_color(tbl,0,RGBColor(0xED,0xE8,0xE0),INK); cell_color(tbl,2,RGBColor(0xE7,0xEE,0xE9),SAGE)
txt(s,9.5,4.5,3.3,0.4,"error medio: 1,4°  ·  imperceptible",11,BRONZE,True,PP_ALIGN.CENTER)
banner(s,"Reproduce los 10 movimientos  ·  y encadena acciones a demanda: correr, saltar y atacar")

# ============================================== 14 GIF
s=slide(); header(s,"V · los resultados","En movimiento: original vs. generado por la red",
                  "la misma patada; a la izquierda el generador, a la derecha la red (en bucle)")
gw,gh=Image.open("ppt/img/resultados.gif").size; Hh=4.5; Wg=Hh*gw/gh
s.shapes.add_picture("ppt/img/resultados.gif",Inches((13.333-Wg)/2),Inches(1.75),Inches(Wg),Inches(Hh))
banner(s,"La red arranca con 16 fotogramas de semilla y genera el resto: el movimiento es casi idéntico")

prs.save("ppt/Proyecto_Ninja_RNN_editable.pptx")
print("guardado:",len(prs.slides._sldIdLst),"diapositivas nativas")
